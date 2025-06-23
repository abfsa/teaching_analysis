import json
import os
from docx import Document
from pptx import Presentation
from lxml import etree
import fitz
from openai import OpenAI
from .util import *

## 设置API

api_key =  os.environ["OPENAI_API_KEY"]
base_url = "https://dashscope.aliyuncs.com/compatible-mode/v1"

client = OpenAI(
  api_key= api_key,
  base_url= base_url
)

def extract_text_docx(file_path):
    doc = Document(file_path)
    full_text = []
    
    # 定义命名空间映射
    nsmap = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

    # 1. 提取普通段落文字
    for para in doc.paragraphs:
        full_text.append(para.text)
    
    # 2. 提取表格中的文字
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                full_text.append(cell.text)
    
    root = etree.fromstring(doc.element.xml)
    txbx_contents = root.xpath('.//w:txbxContent', namespaces=nsmap)
    for txbx in txbx_contents:
        # 查找文本框内的所有段落，并直接提取其中所有文本
        texts = txbx.xpath('.//w:t/text()', namespaces=nsmap)
        if texts:
            full_text.append("".join(texts))
    
    return "\n".join(full_text)

def extract_info_pdf(pdf_path, output_folder="extracted_images"):
    """
    提取 PDF 文档中的所有文字和图片。
    
    参数：
    - pdf_path: PDF 文件路径。
    - output_folder: 存储提取图片的文件夹，默认为 "extracted_images"。
    
    返回：
    - 所有页面的文本拼接成的字符串。
    """
    # 打开 PDF 文件
    doc = fitz.open(pdf_path)
    full_text = []

    # 确保图片存储文件夹存在
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    image_counter = 1

    # 遍历所有页面
    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        # 提取页面文本
        text = page.get_text()
        full_text.append(text)
        
        # 提取页面内的图片（使用 get_images(True) 可以获得更多信息）
        image_list = page.get_images(full=True)
        for img in image_list:
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = os.path.join(output_folder, f"image_page{page_index+1}_{image_counter}.{image_ext}")
            with open(image_filename, "wb") as img_file:
                img_file.write(image_bytes)
            image_counter += 1

    return "\n".join(full_text)

def extract_text_pptx(file_path):
    """从 PPT 文件 (.pptx) 中提取文本"""
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_file(file_path):
    """根据文件后缀自动选择相应的提取方法"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.docx':
        return extract_text_docx(file_path)
    elif ext == '.pptx':
        return extract_text_pptx(file_path)
    elif ext == '.pdf':
        return extract_info_pdf(file_path)
    else:
        raise ValueError("仅支持 .docx , .pdf 和 .pptx 格式文件")

def save_text_to_file(text, output_path):
    """将提取的文本保存到 txt 文件中"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"文本已保存到 {output_path}")

def read_text_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

# 使用 OpenAI 模型提取知识点和关联关系
def extract_knowledge(path, text, style="tree"):
    """
    将文本传递给大语言模型，提取知识点和关联关系，要求返回格式如下：
    {
      "nodes": [
          {"id": "知识点名称", "type": "概念或其他描述"},
          ...
      ],
      "edges": [
          {"source": "知识点1", "target": "知识点2", "relation": "关联类型"},
          ...
      ]
    }
    """
    if style=="tree":
        sample = {
                    "id": "1",
                    "name": "序列决策问题",
                    "content": "",
                    "child": [
                              {"id": "2", "name": "状态", "content": "", "child": []},
                              {"id": "3", "name": "智能体与环境交互", "content": "", "child": []},
                    ]
                }
        attributes = ["重点", "难点", "案例", "总结", "考点", "概述", "实操/训练", "练习", "问题(引例）", "项目/任务/步骤", "外延", "讨论", "情景引入",
    "实验", "岗位", "证书", "比赛"]
        prompt = (f"""请根据以下文本内容，构建一棵树状的知识图谱。要求：
                    1. 顶层节点为本章的主要主题（例如“序列决策问题”、“马尔可夫决策过程”、“值迭代”等）。
                    2. 每个顶层节点下分为若干子节点，表示该主题下的关键概念、定义、模型、方法或例子。
                    3. 每个子节都包含一个属性，属性包括了{attributes}等。
                    4. 子节点下可以继续分层，展示更细粒度的知识点或解释。
                    5. 只需要提取所给文字内容包含的知识点和解释，不必添加额外的知识点和解释。
                    6. 输出格式采用 JSON 树结构，形如：
                    {sample}
                    注意：请严格按照给定的json格式生成，每个节点确保有"id","name","content","child"四个属性。
                    请根据下面的文本内容生成知识图谱：{text}""")

    else:
        sample = {
                        "nodes": [
                        {"id": "1", "name": "序列决策问题"},
                        {"id": "2", "name": "状态"},
                        {"id": "3", "name": "动作"},
                        {"id": "4", "name": "奖励"}
                     ],
                        "edges": [
                        {"source": "1", "relation": "包含", "target": "状态"},
                        {"source": "1", "relation": "包含", "target": "动作"},
                        {"source": "1", "relation": "包含", "target": "奖励"},
                        {"source": "状态", "relation": "影响", "target": "奖励"}
                        ]
                        }
        prompt = (f"""
                    请根据以下文本内容抽取所有核心知识点，并建立它们之间的连接关系，构造一个网状的知识图谱。要求：
                    1. 从文本中识别出所有重要概念、定义、模型、算法、例子等作为节点。
                    2. 针对每对知识点，识别它们之间的关系，如“属于”、“包含”、“导致”、“相互影响”等。
                    3. 输出格式为 JSON 格式，包含两个部分：一个是所有节点列表（每个节点含有唯一ID和名称），另一个是所有关系列表（每个关系包含“起始节点ID”、“关系类型”、“结束节点ID”）。
                    例如：
                        {sample}   

                        请根据下面的文本内容生成知识图谱：{text}""")
    
    response = client.chat.completions.create(
        model="qwen-max",  # 或其他你有权限使用的模型
        messages=[
            {"role": "system", "content": "你是一个知识图谱构建专家。"},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
        #max_tokens=1024
    )
    #print(response)
    result = response.choices[0].message.content.strip()
    with open(os.path.join(path, 'tree1.json'), 'w', encoding = "utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=4)

@retry_on_failure(max_retries=2)
def generate_document_tree(path):
    if path is None:
        return {}
    # 打开文件并按行读取内容
    if os.path.isfile(path):
        try:
            text = extract_text_from_file(path)
        except ValueError:
            raise ValueError(f"不支持的文件类型:{path}")
        return extract_knowledge(os.path.dirname(path), text, style="tree")


    supported_ext = ('.docx', '.pptx', '.pdf')
    all_text = []
    
    for filename in os.listdir(path):
        if filename.lower().endswith(supported_ext):
            file_path = os.path.join(path, filename)
            try:
                text = extract_text_from_file(file_path)
                all_text.append(f"=== 文件: {filename} ===\n{text}\n\n")
            except Exception as e:
                print(f"处理文件 {filename} 时出错: {str(e)}")
    
    output_file = "output.txt"
    if all_text:
        with open(os.path.join(path, output_file), 'w', encoding='utf-8') as f:
            f.write("\n".join(all_text))
        print(f"所有文本已提取并保存到 {output_file}")
    else:
        print("未找到支持的文档文件")
        with open(os.path.join(path, 'tree2.json'), 'w', encoding = "utf-8") as f:
            json.dump({}, f, ensure_ascii=False, indent=4)
            return {}
    file_path = os.path.join(path, 'outline.txt')
    text = read_text_file(file_path)

    response = extract_knowledge(text,"tree")
    return response